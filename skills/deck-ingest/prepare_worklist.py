"""deck-ingest prepare-phase helper.

Deterministic. Extracts per-slide text + speaker notes from .pptx files
via python-pptx and writes the worklist + runbook for Claude Code's
execute phase. No LLM work happens here.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path


def slugify(text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9]+", "-", text.lower())
    return text.strip("-")


def extract_slides(pptx_path: Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise SystemExit(
            "python-pptx is required for deck-ingest — install with: pip install -e '.[ingest]'"
        ) from e

    prs = Presentation(str(pptx_path))
    slides = []
    for n, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets: list[str] = []
        has_images = False
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                has_images = True
            if not shape.has_text_frame:
                continue
            text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                bullets.extend(t for t in text.splitlines() if t.strip())
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        density = len(title) + sum(len(b) for b in bullets) + len(notes)
        slides.append(
            {
                "n": n,
                "title": title,
                "bullets": bullets,
                "notes": notes,
                "has_images": has_images,
                "text_density": density,
            }
        )
    return slides


def prepare(decks_dir: Path, artifacts_dir: Path, content_dir: Path,
            template_path: Path, channel: str) -> None:
    artifacts_decks = artifacts_dir / "decks"
    artifacts_decks.mkdir(parents=True, exist_ok=True)

    decks = sorted(decks_dir.glob("*.pptx"))
    entries = []
    for deck in decks:
        slug = slugify(deck.stem)
        slides = extract_slides(deck)
        slides_path = artifacts_decks / f"{slug}_slides.json"
        slides_path.write_text(json.dumps(slides, indent=2), encoding="utf-8")
        low_text = [s["n"] for s in slides if s["text_density"] < 40 and s["has_images"]]
        entries.append(
            {
                "slug": slug,
                "title": deck.stem.replace("_", " ").replace("-", " ").strip(),
                "slides_path": str(slides_path),
                "source_path": str(deck),
                "channel": channel,
                "output_path": str(content_dir / channel / f"{slug}.md"),
                "slide_count": len(slides),
                "low_text_slides": low_text,
            }
        )

    worklist_path = artifacts_dir / "deck_worklist.json"
    worklist_path.write_text(json.dumps(entries, indent=2), encoding="utf-8")

    runbook = template_path.read_text(encoding="utf-8")
    runbook = (
        runbook.replace("{{worklist_path}}", str(worklist_path))
        .replace("{{total_count}}", str(len(entries)))
        .replace("{{channel}}", channel)
    )
    runbook_path = artifacts_dir / "deck_runbook.md"
    runbook_path.write_text(runbook, encoding="utf-8")

    print(f"Worklist: {len(entries)} deck(s) → {worklist_path}")
    for e in entries:
        note = f", {len(e['low_text_slides'])} low-text slide(s)" if e["low_text_slides"] else ""
        print(f"  - {e['slug']}: {e['slide_count']} slides{note}")
    print(f"Runbook: {runbook_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Prepare the deck ingest runbook")
    parser.add_argument("--decks-dir", type=Path, default=Path("raw/decks"))
    parser.add_argument("--artifacts-dir", type=Path, default=Path("artifacts"))
    parser.add_argument("--content-dir", type=Path, default=Path("content"))
    parser.add_argument("--channel", default="onboarding")
    parser.add_argument(
        "--template", type=Path,
        default=Path("skills/deck-ingest/runbook_template.md"),
    )
    args = parser.parse_args()
    if not args.decks_dir.is_dir():
        raise SystemExit(f"No {args.decks_dir}/ directory — drop .pptx files there first.")
    prepare(args.decks_dir, args.artifacts_dir, args.content_dir, args.template, args.channel)


if __name__ == "__main__":
    main()
