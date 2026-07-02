"""Generate the demo deck for the deck-ingest example.

Creates examples/decks/new-analyst-orientation.pptx — a small, fabricated
orientation deck with speaker notes, used to demonstrate deck-ingest
end-to-end. Requires python-pptx (pip install -e '.[ingest]').
"""

from __future__ import annotations

from pathlib import Path

SLIDES = [
    (
        "New Analyst Orientation",
        ["Model Analytics job family", "Your first week, mapped"],
        "Welcome deck for day one. Everything here also lives in Atlas — this "
        "deck is the guided tour, Atlas is the reference.",
    ),
    (
        "Where everything lives",
        ["Atlas is the front door", "Channels: Onboarding, Models, Compliance, "
         "Lessons, Tooling", "One space per model: whitepaper + user guide + monitoring plan"],
        "Emphasize: if a document is not in Atlas or linked from it, treat it as "
        "not existing. No shared-drive archaeology.",
    ),
    (
        "Your first-week checklist",
        ["Read the Welcome and First 90 days pages", "Skim the glossary",
         "NPI policy BEFORE touching any extract", "Open one model space end to end"],
        "The NPI point is non-negotiable — completing NPI training precedes any "
        "data access request.",
    ),
    (
        "The three documents of every model",
        ["Whitepaper: why it works — methodology, assumptions, limitations",
         "User guide: how to run it — inputs, steps, checks",
         "Monitoring plan: how we know it still works — metrics, thresholds"],
        "Analogy that lands well: theory, procedure, and smoke detector.",
    ),
    (
        "Monitoring in one slide",
        ["Green / Amber / Red per metric", "Amber: document within 10 business days",
         "Red: same-day notification, remediation plan in 30 days",
         "Two consecutive ambers = red"],
        "Numbers here are the CRE PD defaults; each model's plan is authoritative. "
        "The two-ambers rule came out of the 2024 assumption drift lesson.",
    ),
    (
        "When things look wrong",
        ["Never fix silently", "Raise to the model owner and log it",
         "The wrongness is evidence — do not erase it"],
        "Tell the vendor data gap story here: an analyst hand-corrected two scores "
        "and erased the strongest signal of a systemic data issue.",
    ),
    (
        "Who to ask",
        ["Methodology → model owner", "Data and NPI → data steward",
         "Everything else → team channel, tag the lead"],
        "Close with: misrouted questions are cheap, unasked ones are not.",
    ),
]


def main() -> None:
    from pptx import Presentation

    out = Path(__file__).parent / "new-analyst-orientation.pptx"
    prs = Presentation()
    for title, bullets, notes in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "\n".join(bullets)
        slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(out))
    print(f"Wrote {out} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
